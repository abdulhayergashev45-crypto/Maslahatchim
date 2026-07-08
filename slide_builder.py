"""
Sahnalar ro'yxatidan .pptx taqdimot yaratadi.
Brend palitrasi: to'q ko'k (navy), oltin (gold), tiniq moviy (teal), pushti (rose).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x16, 0x1B, 0x31)
PAPER = RGBColor(0xFB, 0xF6, 0xEA)
GOLD = RGBColor(0xF2, 0xB7, 0x05)
TEAL = RGBColor(0x4F, 0xBD, 0xBA)
ROSE = RGBColor(0xE8, 0x63, 0x7C)
INK = RGBColor(0x1B, 0x23, 0x40)
MUTED = RGBColor(0x6B, 0x71, 0x90)

ACCENTS = [TEAL, GOLD, ROSE]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_background(slide, color: RGBColor):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, size, color,
                  bold=False, align=PP_ALIGN.LEFT, font="Calibri", anchor=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def build_pptx(scenes, title: str, output_path: str):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # --- Title slide (dark) ---
    slide = prs.slides.add_slide(blank)
    _set_background(slide, NAVY)
    _add_textbox(slide, Inches(1), Inches(2.7), Inches(11.3), Inches(1.5),
                 title, 40, PAPER, bold=True, align=PP_ALIGN.LEFT, font="Cambria")
    _add_textbox(slide, Inches(1), Inches(3.9), Inches(11.3), Inches(0.6),
                 "Animatsion darslik taqdimoti", 16, TEAL, align=PP_ALIGN.LEFT)

    # --- Content slides ---
    for i, scene in enumerate(scenes):
        accent = ACCENTS[i % len(ACCENTS)]
        slide = prs.slides.add_slide(blank)
        _set_background(slide, PAPER)

        # icon circle
        circle = slide.shapes.add_shape(9, Inches(0.9), Inches(1.0), Inches(2.0), Inches(2.0))  # 9 = OVAL
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.fill.background()
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = scene.get("emoji", "📘")
        run.font.size = Pt(54)

        # scene number
        _add_textbox(slide, Inches(0.9), Inches(0.4), Inches(2.0), Inches(0.5),
                     f"SAHNA {i + 1:02d}", 12, MUTED, font="Courier New")

        # heading
        _add_textbox(slide, Inches(3.4), Inches(1.15), Inches(9.0), Inches(1.0),
                     scene.get("heading", ""), 32, INK, bold=True, font="Cambria")

        # body text
        _add_textbox(slide, Inches(3.4), Inches(2.15), Inches(9.0), Inches(2.0),
                     scene.get("text", ""), 20, INK, font="Calibri")

        # bottom info block (soft tint, no edge stripe)
        block = slide.shapes.add_shape(1, Inches(0.9), Inches(4.6), Inches(11.5), Inches(2.0))  # 1 = RECTANGLE
        block.fill.solid()
        block.fill.fore_color.rgb = RGBColor(0xF3, 0xEF, 0xE2)
        block.line.fill.background()
        block.shadow.inherit = False
        tf2 = block.text_frame
        tf2.word_wrap = True
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf2.margin_left = Inches(0.4)
        tf2.margin_right = Inches(0.4)
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = f"{i + 1} / {len(scenes)} — darslik davomida shu bosqichni eslab qoling"
        run2.font.size = Pt(12)
        run2.font.italic = True
        run2.font.color.rgb = MUTED

    prs.save(output_path)
    return output_path
